"""
Shared pipeline helpers for the Claude Code -> PDF report workflow.

The end-to-end pipeline runs three stages, all orchestrated from the
``run_report.py`` CLI:

    Stage 1  prepare_for_analysis(source_path)
        Extracts dashboard pages to ``temp/slide_N.png`` (or PBI page PNGs)
        and writes ``temp/analysis_request.json``. For PBIP/PBIX it also
        writes ``temp/pbip_context.json`` so Claude can query the live
        Power BI model via the optional ``powerbi-modeling`` MCP.

    Stage 2  trigger_claude_analysis(request_file, context=...)
        Prints a Claude-Code-facing prompt describing what to read and
        what JSON schema to write into ``temp/insights.json``.

    Stage 3  generate_report(...)  (lib.reporting.generate_report)
        Renders the polished A4 PDF report from ``temp/insights.json``.

This module never imports any PPTX-specific code. It only knows about
PDF report generation.
"""

from __future__ import annotations

import json
import os
import re
import time
from pathlib import Path
from typing import Optional


# ---------------------------------------------------------------------------
# Source-file detection
# ---------------------------------------------------------------------------

SUPPORTED_TYPES = {"pdf", "pptx", "pbip", "pbix"}


def detect_file_type(file_path: str) -> str:
    """Detect the input source type from path/extension/folder contents.

    Returns one of: ``"pdf"``, ``"pptx"``, ``"pbip"``, ``"pbix"``.
    """
    p = Path(file_path)
    suffix = p.suffix.lower()

    if suffix == ".pptx":
        return "pptx"
    if suffix == ".pdf":
        return "pdf"
    if suffix == ".pbip":
        return "pbip"
    if suffix == ".pbix":
        return "pbix"
    if p.is_dir() and any(p.glob("*.pbip")):
        return "pbip"

    raise ValueError(
        f"Unsupported source: {file_path}. "
        "Supported formats: .pdf, .pptx, .pbip, .pbix"
    )


# ---------------------------------------------------------------------------
# Power BI Modeling MCP detection (optional deep-analysis path)
# ---------------------------------------------------------------------------

def _load_mcp_server_config() -> Optional[dict]:
    """Return the ``powerbi-modeling`` MCP server config dict, or None.

    Looks in:
      1. ``./.mcp.json`` in the current project
      2. ``~/.claude/mcp-settings.json`` (Claude Code global)

    Tolerates Windows-style unescaped backslashes that Claude Code writes
    natively into its config file.
    """
    def _read_json(path: Path) -> dict:
        text = path.read_text(encoding="utf-8")
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            fixed = re.sub(r'\\(?!["\\/bfnrtu0-9])', r'\\\\', text)
            return json.loads(fixed)

    candidates = [
        Path(".mcp.json"),
        Path.home() / ".claude" / "mcp-settings.json",
    ]
    for path in candidates:
        if path.exists():
            try:
                cfg = _read_json(path)
                server = cfg.get("mcpServers", {}).get("powerbi-modeling")
                if server:
                    return server
            except Exception:
                pass
    return None


def is_pbi_mcp_ready() -> bool:
    """Silent check: True if the Power BI MCP is installed and points to a real exe."""
    server = _load_mcp_server_config()
    return server is not None and Path(server.get("command", "")).exists()


def announce_pbi_mcp_status() -> bool:
    """Verbose status check used by the prepare stage. Never blocks execution."""
    server = _load_mcp_server_config()

    if server is None:
        print()
        print("  " + "!" * 66)
        print("  !! Power BI Modeling MCP is NOT installed.")
        print("  !! Claude will read screenshots only — no live DAX queries.")
        print("  !! To enable exact DAX values:")
        print("  !!     python setup_pbi_mcp.py")
        print("  !! Then restart Claude Code and re-run this command.")
        print("  " + "!" * 66 + "\n")
        return False

    if not Path(server.get("command", "")).exists():
        print()
        print("  " + "!" * 66)
        print("  !! Power BI MCP is configured in .mcp.json but the executable")
        print("  !! was not found at the registered path. Re-run setup:")
        print("  !!     python setup_pbi_mcp.py --force")
        print("  " + "!" * 66 + "\n")
        return False

    print("  OK Power BI Modeling MCP configured — DAX query mode enabled")
    return True


# ---------------------------------------------------------------------------
# Stage 1: extract
# ---------------------------------------------------------------------------

def prepare_for_analysis(source_path: str) -> str:
    """Extract dashboard pages and write ``temp/analysis_request.json``.

    Dispatches to the per-format extractor and returns the absolute path
    of the request file Claude should read.
    """
    file_type = detect_file_type(source_path)

    if file_type == "pdf":
        from lib.extraction.pdf_extractor import prepare_pdf_for_analysis
        return prepare_pdf_for_analysis(source_path)

    if file_type == "pbip":
        announce_pbi_mcp_status()
        from lib.extraction.pbip_extractor import prepare_pbip_for_analysis
        return prepare_pbip_for_analysis(source_path)

    if file_type == "pbix":
        announce_pbi_mcp_status()
        from lib.extraction.pbix_extractor import prepare_pbix_for_analysis
        return prepare_pbix_for_analysis(source_path)

    # PPTX path — supported as legacy input only.
    return _prepare_pptx_for_analysis(source_path)


def _prepare_pptx_for_analysis(source_path: str) -> str:
    """Extract dashboard images from a Power BI PPTX export."""
    import io
    from PIL import Image
    from pptx import Presentation

    print("=" * 70)
    print("PREPARING SLIDES FOR ANALYSIS")
    print("=" * 70)

    prs = Presentation(source_path)
    Path("temp").mkdir(exist_ok=True)

    slides: list[dict] = []
    print(f"\nExtracting {len(prs.slides)} slides... (skipping slide 1 — title page)")

    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue

        title = _extract_slide_title(slide)
        image_path = f"temp/slide_{idx + 1}.png"

        if _extract_slide_image(slide, image_path):
            slides.append({
                "slide_number": idx + 1,
                "title": title,
                "image_path": image_path,
                "slide_type": _classify_slide_type(title),
            })
            print(f"  OK Slide {idx + 1}: {title[:50]}")

    request_file = "temp/analysis_request.json"
    with open(request_file, "w", encoding="utf-8") as f:
        json.dump({
            "source_file": source_path,
            "source_type": "pptx",
            "total_slides": len(slides),
            "slides": slides,
        }, f, indent=2)

    print(f"\nOK Prepared {len(slides)} slides → {request_file}")
    return request_file


def _extract_slide_image(slide, output_path: str) -> bool:
    import io
    from PIL import Image
    for shape in slide.shapes:
        if shape.shape_type == 13:  # picture
            Image.open(io.BytesIO(shape.image.blob)).save(output_path)
            return True
    return False


def _extract_slide_title(slide) -> str:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            title = shape.text.strip()
            return re.sub(r"[\U00010000-\U0010ffff]", "", title).strip()
    return "Untitled Slide"


def _classify_slide_type(title: str) -> str:
    t = title.lower()
    if "trend" in t or "over time" in t:
        return "trends"
    if "leaderboard" in t or "top" in t:
        return "leaderboard"
    if "health" in t or "overview" in t:
        return "health_check"
    if "habit" in t or "frequency" in t:
        return "habit_formation"
    if "license" in t or "priority" in t:
        return "license_priority"
    return "general"


# ---------------------------------------------------------------------------
# Stage 2: instruct Claude Code to analyse
# ---------------------------------------------------------------------------

def trigger_claude_analysis(request_file: str, context: Optional[str] = None) -> None:
    """Print the Claude Code analysis prompt for the given request.

    Picks the PBIP/MCP variant when the source is PBIP/PBIX and the MCP
    is installed; otherwise falls back to the image-based variant.
    """
    print("\n" + "=" * 70)
    print("STEP 2: CLAUDE CODE ANALYSIS")
    print("=" * 70)

    with open(request_file, "r", encoding="utf-8") as f:
        request = json.load(f)

    is_pbip = (
        request.get("source_type") in ("pbip", "pbix")
        or Path("temp/pbip_context.json").exists()
    )

    if is_pbip and is_pbi_mcp_ready():
        _print_pbip_prompt(request, context=context)
    elif is_pbip:
        _print_image_prompt(request, mcp_missing=True, context=context)
    else:
        _print_image_prompt(request, context=context)


def _print_image_prompt(request: dict, *, mcp_missing: bool = False,
                        context: Optional[str] = None) -> None:
    if mcp_missing:
        print("\n!! MCP not installed — running in IMAGE-ONLY mode. Numbers will be")
        print("!! read visually from screenshots, not queried from the live model.")
        print("!! Install with:  python setup_pbi_mcp.py\n")

    if context:
        print(f"ANALYSIS FOCUS: {context}\n")

    print(f"Claude Code: please analyse {request['total_slides']} dashboard pages.\n")
    print("Pages to analyse:")
    for slide in request["slides"]:
        print(f"  - Page {slide['slide_number']}: {slide['title']}")
        print(f"    Image: {slide['image_path']}")

    print("\n" + "-" * 70)
    print("CLAUDE CODE TASK")
    print("-" * 70)
    print(_IMAGE_PROMPT_BODY)


def _print_pbip_prompt(request: dict, *, context: Optional[str] = None) -> None:
    print(f"\nClaude Code: please analyse this Power BI report "
          f"({request['total_slides']} pages) using the live model via MCP.\n")
    print("Pages:")
    for slide in request["slides"]:
        print(f"  - Page {slide['slide_number']}: {slide['title']} ({slide['slide_type']})")

    print("\n" + "-" * 70)
    print("CLAUDE CODE TASK (PBIP / MCP MODE)")
    print("-" * 70)

    if context:
        print(f"\nANALYSIS FOCUS: {context}\n")

    print(_PBIP_PROMPT_BODY)


_IMAGE_PROMPT_BODY = """
Act as a senior analyst advising an IT decision maker.

For EACH dashboard page above:
  1. Read the screenshot file
  2. Extract every visible number, percentage, and label with EXACT units
     (write "13K" not "13,000"; write "5.7×" not "570%")
  3. Generate a memorable headline that answers "so what?" — not a data dump
  4. Generate up to 3 insights, each formatted "Bold line || Supporting evidence"
  5. Use opportunity framing — never crisis or failure language
  6. Only mention platforms / teams / features visible on THIS page
  7. If a page has no numbers, set headline to "Insufficient data for analysis"
     and leave numbers_used empty

Then synthesise across ALL pages:
  - deck_title    — compelling 5–10 word story-driven title
  - deck_subtitle — "[Platform] · [Org] · [Date range]"
  - executive_summary — 5 bullets, highest business impact first
  - recommendations   — 3–5 specific actions traceable to dashboard data

Save the result to:  temp/insights.json

CHART SPECS (recommended, especially for KPIs and tables)
  Each insight may include a "chart" object that the PDF renderer will
  draw as a vector chart. Set "chart": null when the screenshot already
  tells the story.

  Common chart types: kpi, kpi_row, bar, column, line, area, donut, pie,
  table, treemap, funnel, gauge, heatmap, scatter, waterfall, combo.
  TABLES MUST stay tables — never convert "table" specs to bar/column.

  Examples:
    {"type": "kpi", "value": "4,381", "label": "Active Users"}
    {"type": "kpi_row", "items": [
        {"value": "256", "label": "Active Users"},
        {"value": "33.77", "label": "Weekly Actions"}
    ]}
    {"type": "bar", "title": "Sessions by Org", "highlight": "Finance",
     "data": [{"label": "Finance", "value": 5.5},
              {"label": "Legal",   "value": 3.26}]}
    {"type": "table", "columns": ["Manager", "Users", "Weekly Actions"],
     "rows": [["Dana Bourque", "4", "48.95"],
              ["Matt Sheard",  "52", "38.09"]]}

OUTPUT SCHEMA (temp/insights.json)
  {
    "deck_title": "Compelling story-driven title",
    "deck_subtitle": "Agents · Chat · M365 Copilot · Mar – Jun 2025",
    "executive_summary": ["finding → implication", ...],
    "recommendations":   ["action → expected outcome", ...],
    "slides": [
      {
        "slide_number": 1,
        "title": "Insight-led title (can differ from source)",
        "headline": "Memorable takeaway answering 'so what?'",
        "insights": [
          {"text": "Bold line || Supporting evidence with data",
           "chart": {"type": "kpi", "value": "...", "label": "..."}},
          {"text": "Second insight || Detail",       "chart": null},
          {"text": "Third insight (action) || Detail", "chart": null}
        ],
        "numbers_used": ["123", "45%"]
      }
    ]
  }

When done, run:  python run_report.py --build --input "<source>"
"""


_PBIP_PROMPT_BODY = """
Act as a senior analyst advising an IT decision maker.

This is a Power BI project — the LIVE model is reachable via the
'powerbi-modeling' MCP server. DO NOT estimate numbers from screenshots
when DAX is available. Every number in insights MUST come from a query.

STEP 1 — Read context
    Read temp/pbip_context.json
    (Contains pages, visuals, measures with full DAX, pre-built dax_queries.)

STEP 2 — For each page, execute its DAX queries via the MCP
    Use mcp__powerbi-modeling__dax_query_operations.execute_query.
    Power BI Desktop must be open with this project for the MCP to connect.
    Run every query in pbip_context.json -> dax_queries[n].queries[m].dax.
    The returned table rows ARE your data source — use exact values.

STEP 3 — Drill or filter as needed
    Wrap with CALCULATE / CALCULATETABLE for segment analysis.
    Use DATESYTD / DATESINPERIOD for time-filtered views.
    Compare filtered vs unfiltered for delta / trend insights.

STEP 4 — Translate visuals to chart specs
    PBI visual type        ->  ChartSpec type
    -------------------------------------------------
    tableEx / matrix       ->  "table"     (NEVER collapse to bar/column)
    card / multiRowCard    ->  "kpi" / "kpi_row"
    bar / column / line    ->  "bar" / "column" / "line"
    donut / pie            ->  "donut" / "pie"
    waterfall              ->  "waterfall"
    treemap                ->  "treemap"
    combo / dual-axis      ->  "combo"

STEP 5 — Generate insights
    Same formula — headline + up to 3 "Bold || Detail" insights per page.
    Cite measure names alongside values for traceability.
    Include numbers_used so the PDF appendix can list every cited figure.

STEP 6 — Write temp/insights.json
    Schema (identical to image-mode):
      {
        "deck_title": "...",
        "deck_subtitle": "...",
        "executive_summary": [...],
        "recommendations":   [...],
        "slides": [
          {"slide_number": 1, "title": "...", "headline": "...",
           "insights": [{"text": "Bold || Detail",
                         "chart": {"type": "kpi", "value": "4,381", "label": "..."}}],
           "numbers_used": ["..."]}
        ]
      }

When done, run:  python run_report.py --build --input "<source>"
"""


# ---------------------------------------------------------------------------
# Stage 2.5: verify insights.json before render
# ---------------------------------------------------------------------------

_VANILLA_HEADLINE_PATTERNS = (
    "groups form", "org structure", "overview of", "shows the",
    "displays the", "presents the", "summary of", "breakdown of",
    "distribution of", "there are", "this page shows", "this slide shows",
)


def verify_insights(insights_file: str = "temp/insights.json",
                    request_file: str = "temp/analysis_request.json") -> dict:
    """Validate insights.json for missing charts, weak headlines, gaps.

    Returns ``{"passed": bool, "errors": [...], "warnings": [...]}``.
    Errors block the render; warnings do not.
    """
    warnings: list[str] = []
    errors: list[str] = []

    try:
        with open(insights_file, "r", encoding="utf-8") as f:
            insights = json.load(f)
    except FileNotFoundError:
        return {"passed": False, "warnings": [],
                "errors": [f"Insights file not found: {insights_file}"]}
    except json.JSONDecodeError as e:
        return {"passed": False, "warnings": [],
                "errors": [f"Invalid JSON in {insights_file}: {e}"]}

    slides = insights.get("slides", [])
    if not slides:
        return {"passed": False, "warnings": [],
                "errors": ["No slides found in insights JSON"]}

    # Missing charts
    for s in slides:
        slide_num = s.get("slide_number", "?")
        items = s.get("insights", [])
        has_chart = any(isinstance(it, dict) and it.get("chart") for it in items)
        headline = (s.get("headline") or "").lower().strip()
        if not has_chart and headline != "insufficient data for analysis":
            warnings.append(
                f"Page {slide_num}: no chart spec — PDF will rely on screenshot only."
            )

    # Slide-count check
    try:
        with open(request_file, "r", encoding="utf-8") as f:
            request = json.load(f)
        expected = request.get("total_slides", 0)
        if expected and len(slides) < expected:
            warnings.append(
                f"Slide count mismatch: expected {expected}, got {len(slides)}. "
                "Some pages may be missing from insights."
            )
    except (OSError, json.JSONDecodeError):
        pass

    # Vanilla headlines
    for s in slides:
        headline = (s.get("headline") or "").lower()
        for pat in _VANILLA_HEADLINE_PATTERNS:
            if pat in headline:
                warnings.append(
                    f"Page {s.get('slide_number', '?')}: headline may be too generic "
                    f"(matched '{pat}'). Headlines should answer 'so what?' for an "
                    "executive — not describe the chart."
                )
                break

    # Required deck-level fields
    if not insights.get("executive_summary"):
        errors.append("Missing 'executive_summary' — required for the summary page.")
    if not insights.get("recommendations"):
        warnings.append("Missing 'recommendations' — strongly recommended.")
    title = insights.get("deck_title")
    if not title or title == "Executive Insights":
        warnings.append("Deck title is generic. Generate a compelling, specific title.")

    passed = len(errors) == 0

    print("\n" + "=" * 70)
    print("INSIGHT VERIFICATION REPORT")
    print("=" * 70)
    for e in errors:
        print(f"  ERROR  {e}")
    for w in warnings:
        print(f"  WARN   {w}")
    if not errors and not warnings:
        print("  OK  All checks passed — insights look good")
    print(f"\n  Result: {'PASS' if passed else 'FAIL'} "
          f"({len(errors)} error(s), {len(warnings)} warning(s))")
    print("=" * 70)

    return {"passed": passed, "warnings": warnings, "errors": errors}


# ---------------------------------------------------------------------------
# Stage 2 → Stage 3 bridge: poll for insights.json
# ---------------------------------------------------------------------------

def wait_for_insights(insights_path: str = "temp/insights.json",
                      max_wait: int = 300, interval: int = 2) -> bool:
    """Poll for the analyst (Claude) to finish writing insights.json.

    Returns True as soon as the file exists and parses with at least one
    slide. Returns False if it never appears within ``max_wait`` seconds.
    """
    print("\n" + "=" * 70)
    print("Waiting for temp/insights.json...")
    print("=" * 70)

    p = Path(insights_path)
    elapsed = 0
    while elapsed < max_wait:
        if p.exists():
            try:
                data = json.loads(p.read_text(encoding="utf-8"))
                if data.get("slides"):
                    print(f"OK insights ready ({elapsed}s)")
                    return True
            except (json.JSONDecodeError, KeyError):
                pass
        time.sleep(interval)
        elapsed += interval

    print(f"WARN insights file not ready after {max_wait}s.")
    return False


# ---------------------------------------------------------------------------
# Cleanup helpers
# ---------------------------------------------------------------------------

def cleanup_insight_artefacts() -> None:
    """Remove per-run insight artefacts after a successful build."""
    for fname in ("temp/insights.json", "temp/write_insights.py"):
        try:
            os.remove(fname)
        except FileNotFoundError:
            pass
